import os
import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

pptx_path = r"c:\Users\asus\OneDrive - Trường Ngôi Sao Hoàng Mai\Desktop\phiếu nhận xét học sinh cuối năm\Phieu_bao_cao_ket_qua_hoc_tap_ca_nam_lop_4C.pptx"

if not os.path.exists(pptx_path):
    print("Cannot find PowerPoint file!")
    sys.exit(1)

prs = Presentation(pptx_path)
print(f"Total slides in PPTX: {len(prs.slides)}")

for idx, slide in enumerate(prs.slides):
    slide_text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text.append(run.text)
    
    text_content = " ".join(slide_text)
    # Search for "Học sinh:" or "Họ và tên:" or just print some matching names
    print(f"Slide {idx}: {text_content[:200]}")
