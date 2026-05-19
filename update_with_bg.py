import os
from pptx import Presentation
import json
import shutil

base_dir = r"c:\Users\asus\OneDrive - Trường Ngôi Sao Hoàng Mai\Desktop\phiếu nhận xét học sinh cuối năm\web_so_lo_xo"
pptx_path = r"c:\Users\asus\OneDrive - Trường Ngôi Sao Hoàng Mai\Desktop\phiếu nhận xét học sinh cuối năm\Pink and White Playful  Illustrative Back to School Poster (A4).pptx"
out_dir = os.path.join(base_dir, 'assets', 'students')
os.makedirs(out_dir, exist_ok=True)

prs = Presentation(pptx_path)

student_images = []
img_count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.shape_type == 13: # Picture
            if abs(shape.left - 693886) < 200000 and abs(shape.top - 3655565) < 200000:
                image_bytes = shape.image.blob
                ext = shape.image.ext
                filename = f"student_{img_count}.{ext}"
                img_path = os.path.join(out_dir, filename)
                with open(img_path, 'wb') as f:
                    f.write(image_bytes)
                student_images.append(f"assets/students/{filename}")
                img_count += 1
                break # Only one per slide

print(f"Extracted {img_count} student images")

# Copy back.png
shutil.copy2(r"c:\Users\asus\OneDrive - Trường Ngôi Sao Hoàng Mai\Desktop\phiếu nhận xét học sinh cuối năm\back.png", os.path.join(base_dir, 'assets', 'back.png'))

# Re-read students
with open(os.path.join(base_dir, 'names.json'), 'r', encoding='utf-8') as f:
    data = json.load(f)

students = []
for row in data:
    try:
        stt = row.get("Unnamed: 0")
        name = row.get("Unnamed: 1")
        if isinstance(stt, (int, float)) and name and isinstance(name, str) and not "HỌ VÀ TÊN" in name.upper():
            students.append(name)
    except:
        pass

# Generate HTML
cards_html = []
for i, name in enumerate(students):
    img = student_images[i % len(student_images)] if student_images else 'assets/boy.png'
    card = f'''
        <div class="student-card">
            <img src="{img}" alt="Học sinh">
            <h4>{name}</h4>
        </div>'''
    cards_html.append(card)

cards_str = '\n'.join(cards_html)

index_html = f'''<!DOCTYPE html>
<html lang="vi">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Kỷ Yếu Lớp 4C</title>
    <link rel="stylesheet" href="style.css">
    <link href="https://fonts.googleapis.com/css2?family=Edu+VIC+WA+NT+Beginner:wght@400..700&display=swap" rel="stylesheet">
</head>
<body>
    <div class="mobile-notebook">
        <div class="content-area">
            <h1>Trường Tiểu học Mỹ Thuận 1</h1>
            <div class="teacher-box">
                <h2>GVCN: Cô Nguyễn Thị Thu Cúc</h2>
                <p>Kỷ yếu cuối năm học 2023-2024</p>
            </div>
            
            <h2 class="class-title">Danh sách lớp 4C</h2>
            <div class="student-grid">
{cards_str}
            </div>
        </div>
    </div>
</body>
</html>
'''
with open(os.path.join(base_dir, 'index.html'), 'w', encoding='utf-8') as f:
    f.write(index_html)

style_css = '''
:root {
    --font-main: 'Edu VIC WA NT Beginner', cursive;
}

body {
    margin: 0;
    padding: 0;
    background-color: #ff9fb2;
    display: flex;
    justify-content: center;
    align-items: center;
    min-height: 100vh;
    font-family: var(--font-main);
}

.mobile-notebook {
    width: 100%;
    max-width: 600px; /* like a mobile view */
    height: 100vh;
    max-height: 1000px;
    background-image: url('assets/back.png');
    background-size: cover;
    background-position: center;
    position: relative;
    box-shadow: 0 10px 30px rgba(0,0,0,0.2);
}

/* The actual lined area of the notebook in the image */
.content-area {
    position: absolute;
    /* Adjust these based on the actual lines in the image */
    top: 22%;
    bottom: 25%;
    left: 42%; /* Right side of the spiral */
    right: 12%; /* Before the tabs */
    overflow-y: auto;
    padding: 10px;
}

.content-area::-webkit-scrollbar { width: 5px; }
.content-area::-webkit-scrollbar-thumb { background: rgba(255,107,129,0.5); border-radius: 5px; }

h1 {
    color: #ff4757;
    font-size: 1.8rem;
    text-align: center;
    margin-top: 0;
    margin-bottom: 10px;
}

.teacher-box {
    text-align: center;
    border-bottom: 2px dashed #ff4757;
    padding-bottom: 10px;
    margin-bottom: 15px;
}
.teacher-box h2 {
    font-size: 1.2rem;
    color: #2f3542;
    margin: 0;
}
.teacher-box p {
    margin: 5px 0 0 0;
    font-size: 1rem;
    color: #57606f;
}

.class-title {
    text-align: center;
    font-size: 1.5rem;
    color: #2f3542;
    margin-bottom: 15px;
}

.student-grid {
    display: grid;
    grid-template-columns: repeat(2, 1fr);
    gap: 15px;
    padding-bottom: 20px;
}

.student-card {
    text-align: center;
    background: rgba(255, 255, 255, 0.6);
    padding: 5px;
    border-radius: 8px;
    box-shadow: 0 2px 5px rgba(0,0,0,0.1);
}

.student-card img {
    width: 60px;
    height: 60px;
    border-radius: 50%;
    object-fit: cover;
    border: 2px solid #ff4757;
    margin-bottom: 5px;
}

.student-card h4 {
    margin: 0;
    font-size: 0.9rem;
    color: #2d3436;
}
'''
with open(os.path.join(base_dir, 'style.css'), 'w', encoding='utf-8') as f:
    f.write(style_css)

print("Setup complete")
