import shutil
import os
from pptx import Presentation

src = r"c:\Users\asus\OneDrive - Trường Ngôi Sao Hoàng Mai\Desktop\phiếu nhận xét học sinh cuối năm\Pink and White Playful  Illustrative Back to School Poster (A4).pptx"
dst = r"C:\Users\asus\AppData\Local\Temp\poster.pptx"

shutil.copy(src, dst)
prs = Presentation(dst)
print(f"Number of slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    if i > 2: break
    print(f"Slide {i+1}:")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print("  Text:", shape.text[:50].replace('\n', ' '))
