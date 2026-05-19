import collections 
import collections.abc
from pptx import Presentation
import os

prs = Presentation(r"c:\Users\asus\OneDrive - Trường Ngôi Sao Hoàng Mai\Desktop\phiếu nhận xét học sinh cuối năm\Pink and White Playful  Illustrative Back to School Poster (A4).pptx")

print(f"Number of slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for shape in slide.shapes:
        if hasattr(shape, 'image'): # Check if it's a picture or has an image
            print(f"  Image: name='{shape.name}', left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")
        elif shape.shape_type == 14: # PLACEHOLDER
            print(f"  Placeholder: name='{shape.name}'")
    if i > 4:
        break
